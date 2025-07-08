from pathlib import Path
import pymupdf
import docx2txt
from pptx import Presentation
from backend.tokenizer import tokenize

def match_extractor(path):
    ext = Path(path).suffix.lower().lstrip(".")

    extractors = {
        "pdf": pdf,
        "txt": txt,
        "docx": docx,
        "doc": docx,
        "md": markdown,
        "pptx": pptx
    }

    if ext not in extractors:
        print(f"Unsupported file type: {ext}")
        return None
    else:
        return extractors[ext]

def txt(file_path, tokenize: bool = False):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        raw_content = f.read()
    if tokenize:
        return [tokenize(raw_content)]
    return [raw_content]

def pdf(file_path, tokenize: bool = False):
    pages = []
    doc = pymupdf.open(file_path)
    for page in doc:
        text = page.get_text()
        if tokenize:
            pages.append(tokenize(text))
        else:
            pages.append(text)
    return pages

def docx(file_path, tokenize: bool = False):
    raw_content = docx2txt.process(file_path)
    if tokenize:
        return [tokenize(raw_content)]
    return [raw_content]

def pptx(file_path, tokenize: bool = False):
    prs = Presentation(file_path)
    slides = []
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text+"\n"
        if tokenize:
            slides.append(tokenize(slide_text))
        else:
            slides.append(slide_text)
    return slides

def markdown(file_path, tokenize: bool = False):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        raw_content = f.read()
    if tokenize:
        return [tokenize(raw_content)]
    return [raw_content]