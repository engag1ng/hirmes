"""
Document reading utilities for extracting text.

Typical usage:
    file_func = matchextractor(file_path)
    content = file_func(file_path, False)
"""

from pathlib import Path
import pymupdf
import docx2txt
from pptx import Presentation

def match_extractor(path: str):
    """Higher-level function returning appropriate extraction function for file.

    Args:
        path: String of full path to file.

    Returns:
        func: Function for extracting content from file.
    """

    ext = Path(path).suffix.lower().lstrip(".")

    extractors = {
        "pdf": pdf,
        "txt": txt,
        "docx": docx,
        "md": markdown,
        "pptx": pptx
    }

    if ext in extractors:
        return extractors[ext]
    return None

def txt(path: str) -> list | None:
    """Extracts text from txt file.

    Args:
        path: String of full path to file.
        
    Returns:
        list: One item list of text string.
        None: If extraction fails.
    """

    if not Path(path).is_file():
        return None
    with open(path, 'r', encoding='utf-8', errors='ignore') as file:
        raw_content = file.read()
    return [raw_content]

def pdf(path: str) -> list | None:
    """Extracts text from pdf file.

    Args:
        path: String of full path to file.

    Returns:
        list: One string per page.
        None: If extraction fails.
    """

    try:
        document = pymupdf.open(path)
        pages = [page.get_text() for page in document]
        return pages
    except Exception:
        return None

def docx(path: str) -> list | None:
    """Extracts text from docx file.

    Args:
       path: String of full path to file. 
    
    Returns:
        list: One item list of text string.
        None: If extraction fails.
    """

    try:
        raw_content = docx2txt.process(path)
        return [raw_content]
    except Exception:
        return None

def pptx(path: str) -> list | None:
    """Extracts text from pptx file.

    Args:
        path: String of full path to file.

    Returns:
        list: One string per page.
        None: If extraction fails.
    """
    try:
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text+"\n"
                slides.append(slide_text)
        return slides
    except Exception:
        return None

def markdown(path: str) -> list | None:
    """Extracts text from markdown file.

    Args:
        path: String of full path to file.

    Returns:
        list: One item list of text string. 
        None: If extraction fails.
    """

    if not Path(path).is_file():
        return None
    with open(path, 'r', encoding='utf-8', errors='ignore') as file:
        raw_content = file.read()
    return [raw_content]
