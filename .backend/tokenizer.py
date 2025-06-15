import pymupdf4llm
import docx2txt
from pptx import Presentation
import re

## File types
def txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        raw_content = f.read()
    return tokenize(raw_content)

def pdf(file_path):
    return tokenize(pymupdf4llm.to_markdown(file_path))

def docx(file_path):
    return tokenize(docx2txt.process(file_path))

def pptx(file_path):
    full_string = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_string += shape.text+"\n"
    return tokenize(full_string)

def markdown(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        raw_content = f.read()
    return tokenize(raw_content)


## Tokenization   
def tokenize(content):
    symbols_to_remove = ['%', '^', '&', '*', '~', '[', ']']
    table = str.maketrans('', '', ''.join(symbols_to_remove))
    cleaned_text = content.translate(table)

    # More comprehensive URL pattern
    url_pattern = r'\b(?:https?://)?(?:www\.)?[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}(?:/[^\s]*)?'
    urls = re.findall(url_pattern, cleaned_text)
    no_urls = re.sub(url_pattern, '', cleaned_text)
    clean_urls = []
    for url in urls:
        if '(' in url and url.count('http') > 1:
            url = url.split('(')[0]
        url = url.rstrip(').,;')
        if not re.search(r'\.\w{2,5}(/|$)', url):  # ensures domain or file extension
            continue

        clean_urls.append(url)
    clean_urls = list(set(clean_urls))

    symbols_to_remove2 = ['(', ')', ':', '/']
    table2 = str.maketrans('', '', ''.join(symbols_to_remove2))
    second_clean = no_urls.translate(table2)

    date_pattern = r'\b(0?[1-9]|[12][0-9]|3[01])\.(0?[1-9]|1[0-2])\.(\d{4})\b'
    dates = ['.'.join(date) for date in re.findall(date_pattern, second_clean)]
    no_dates = re.sub(date_pattern, '', second_clean)

    split_by_line = [line.lstrip('#') for line in no_dates.split("\n") if line != "" and line != " "]

    filter_set = set(['“', '-', '_', '.', ','] + [str(i) for i in range(10)] + [f'{i:02}' for i in range(10)])

    split_by_word = []

    for line in split_by_line:
        for text in line.split():
            cleaned = text.strip('",.“”>`!?')

            if len(cleaned) == 1 and cleaned in filter_set:
                continue

            for sym in filter_set:
                cleaned = re.sub(f"{re.escape(sym)}{{2,}}", '', cleaned)

            if cleaned:
                split_by_word += re.split(r"[’']+", cleaned.lower())

    stop_list_filtered = [token for token in split_by_word if token not in stop_list()]

    filtered_complete = stop_list_filtered + dates + clean_urls

    freq_dict = {}
    for token in filtered_complete:
        if token in freq_dict:
            freq_dict[token] += 1
        else:
            freq_dict[token] = 1

    return sorted(freq_dict.items(), key=lambda x: x[1])

def stop_list():
    with open(".backend/stop_list.txt", 'r', encoding="utf-8") as file:
        return file.read().split("\n")

print(pdf("testpdf.pdf"))